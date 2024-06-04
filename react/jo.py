from pptx import Presentation
from pptx.util import Inches

# Crear una nueva presentación
prs = Presentation()

# Diapositiva 1: Título
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Crecimiento Poblacional de Perú (2010-2022)"
subtitle.text = "Análisis de la población de Perú en los últimos 12 años"

# Diapositiva 2: Introducción
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Introducción"
content.text = (
    "Esta presentación analiza el crecimiento poblacional de Perú entre los años 2010 y 2022. "
    "Los datos presentados muestran cómo ha cambiado la población durante esta década."
)

# Diapositiva 3: Datos
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_3.shapes.title
title.text = "Datos de Población (2010-2022)"
rows = len(years_updated) + 1
cols = 2
left = Inches(2.0)
top = Inches(1.5)
width = Inches(6.0)
height = Inches(3.0)
table = slide_3.shapes.add_table(rows, cols, left, top, width, height).table

# Configurar el encabezado
table.cell(0, 0).text = 'Año'
table.cell(0, 1).text = 'Población (millones)'

# Llenar la tabla con datos
for i, year in enumerate(years_updated):
    table.cell(i+1, 0).text = str(year)
    table.cell(i+1, 1).text = str(population_updated[i])

# Diapositiva 4: Gráfica
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])
title = slide_4.shapes.title
title.text = "Gráfica del Crecimiento Poblacional"
left = Inches(1.0)
top = Inches(1.5)
height = Inches(4.5)
pic = slide_4.shapes.add_picture("/mnt/data/Peru_Population_Growth_Updated.png", left, top, height=height)

# Diapositiva 5: Análisis
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Análisis"
content.text = (
    "La gráfica muestra una tendencia de crecimiento constante en la población de Perú "
    "durante el período 2010-2022. Se observa un aumento gradual cada año, lo que indica "
    "un crecimiento sostenido y saludable de la población."
)

# Diapositiva 6: Conclusión
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]
title.text = "Conclusión"
content.text = (
    "El análisis del crecimiento poblacional de Perú durante la última década revela una tendencia positiva. "
    "Este crecimiento continuo puede ser atribuido a factores como mejoras en la salud, economía y migración. "
    "Es fundamental seguir monitoreando estos datos para planificar adecuadamente el futuro del país."
)

# Guardar la presentación
prs.save("/mnt/data/Crecimiento_Poblacional_Peru_Actualizado.pptx")
